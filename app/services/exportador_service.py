"""Infraestructura de exportación GENÉRICA y transversal: DOCX · PDF · XLSX.

Complementa a export_service (que es específico del curso). Cualquier módulo arma su contenido
y llama a `exportar(formato, payload)`. Formato por contenido (lo más riguroso):
  · DOCX  → documentos narrativos (manuscrito, informe): Word editable.
  · PDF   → entregable de layout fijo (impresión del informe/manuscrito).
  · XLSX  → datos tabulares (corpus, efectos, ítems, referencias, notas).

Documento narrativo (docx/pdf):
  {"titulo": str, "secciones": [{"heading": str|None, "nivel": 1|2, "texto": str}],
   "tablas": [{"titulo": str, "headers": [...], "rows": [[...]]}]}
Libro tabular (xlsx):
  {"hojas": [{"nombre": str, "headers": [...], "rows": [[...]]}]}

El texto admite **negrita** con markdown simple. Dependencias: reportlab, openpyxl, python-docx.
"""
from __future__ import annotations

import io
import re


def _split_bold(texto: str):
    partes, i = [], 0
    for m in re.finditer(r"\*\*(.+?)\*\*", texto or ""):
        if m.start() > i:
            partes.append((texto[i:m.start()], False))
        partes.append((m.group(1), True))
        i = m.end()
    if i < len(texto or ""):
        partes.append((texto[i:], False))
    return partes or [("", False)]


# ───────────────────────────── XLSX (openpyxl)
def to_xlsx(doc) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment

    # Acepta el payload completo ({hojas, imagenes}) o una lista de hojas (compatibilidad).
    hojas = doc.get("hojas") if isinstance(doc, dict) else doc
    imagenes = doc.get("imagenes") if isinstance(doc, dict) else None

    wb = Workbook()
    wb.remove(wb.active)
    hdr_font = Font(bold=True, color="FFFFFF")
    hdr_fill = PatternFill("solid", fgColor="4B3F72")
    for h in (hojas or [{"nombre": "Datos", "headers": [], "rows": []}]):
        ws = wb.create_sheet((h.get("nombre") or "Hoja")[:31])
        headers = h.get("headers") or []
        if headers:
            ws.append([str(x) for x in headers])
            for c in ws[1]:
                c.font = hdr_font
                c.fill = hdr_fill
                c.alignment = Alignment(vertical="center")
        for row in (h.get("rows") or []):
            ws.append(["" if v is None else v for v in row])
        for col in ws.columns:
            largo = max((len(str(c.value)) for c in col if c.value is not None), default=10)
            ws.column_dimensions[col[0].column_letter].width = min(max(largo + 2, 10), 60)
        if headers:
            ws.freeze_panes = "A2"
    # Gráficos (PNG) en una hoja aparte.
    for im in (imagenes or []):
        if not im.get("png"):
            continue
        try:
            from openpyxl.drawing.image import Image as XLImage
            ws = wb.create_sheet((im.get("titulo") or "Gráfico")[:31])
            pic = XLImage(io.BytesIO(im["png"]))
            ws.add_image(pic, "B2")
        except Exception:
            pass
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ───────────────────────────── DOCX (python-docx)
def to_docx(doc: dict) -> bytes:
    from docx import Document
    from docx.shared import Pt

    d = Document()
    if doc.get("titulo"):
        d.add_heading(doc["titulo"], level=0)
    for sec in (doc.get("secciones") or []):
        if sec.get("heading"):
            d.add_heading(sec["heading"], level=int(sec.get("nivel", 1)))
        if sec.get("texto"):
            for parrafo in str(sec["texto"]).split("\n"):
                p = d.add_paragraph()
                for seg, bold in _split_bold(parrafo):
                    run = p.add_run(seg)
                    run.bold = bold
                    run.font.size = Pt(11)
    for im in (doc.get("imagenes") or []):
        if not im.get("png"):
            continue
        if im.get("titulo"):
            d.add_heading(im["titulo"], level=2)
        try:
            from docx.shared import Inches
            d.add_picture(io.BytesIO(im["png"]), width=Inches(6))
        except Exception:
            pass
    for t in (doc.get("tablas") or []):
        if t.get("titulo"):
            d.add_heading(t["titulo"], level=2)
        headers = t.get("headers") or []
        tbl = d.add_table(rows=1, cols=max(1, len(headers)))
        try:
            tbl.style = "Light Grid Accent 1"
        except Exception:
            pass
        for j, hcell in enumerate(headers):
            tbl.rows[0].cells[j].paragraphs[0].add_run(str(hcell)).bold = True
        for row in (t.get("rows") or []):
            cells = tbl.add_row().cells
            for j, v in enumerate(row[:len(headers)] if headers else row):
                cells[j].text = "" if v is None else str(v)
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


# ───────────────────────────── PDF (reportlab)
def to_pdf(doc: dict) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle

    ss = getSampleStyleSheet()
    h0 = ParagraphStyle("H0", parent=ss["Title"], fontSize=16, leading=20, spaceAfter=10)
    h1 = ParagraphStyle("H1", parent=ss["Heading2"], fontSize=12,
                        textColor=colors.HexColor("#4B3F72"), spaceBefore=10, spaceAfter=4)
    body = ParagraphStyle("B", parent=ss["BodyText"], fontSize=10, leading=15, spaceAfter=6, alignment=4)

    def esc(s):
        s = (s or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        return re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", s)

    story = []
    if doc.get("titulo"):
        story.append(Paragraph(esc(doc["titulo"]), h0))
    for sec in (doc.get("secciones") or []):
        if sec.get("heading"):
            story.append(Paragraph(esc(sec["heading"]), h1))
        if sec.get("texto"):
            for parrafo in str(sec["texto"]).split("\n"):
                if parrafo.strip():
                    story.append(Paragraph(esc(parrafo), body))
    for im in (doc.get("imagenes") or []):
        if not im.get("png"):
            continue
        if im.get("titulo"):
            story.append(Paragraph(esc(im["titulo"]), h1))
        try:
            from reportlab.platypus import Image as RLImage
            iw = 15 * cm
            story.append(Spacer(1, 4))
            story.append(RLImage(io.BytesIO(im["png"]), width=iw, height=iw * 340.0 / 640.0))
        except Exception:
            pass
    for t in (doc.get("tablas") or []):
        if t.get("titulo"):
            story.append(Paragraph(esc(t["titulo"]), h1))
        headers = [str(x) for x in (t.get("headers") or [])]
        data = [headers] + [["" if v is None else str(v) for v in row] for row in (t.get("rows") or [])]
        if len(data) > 1:
            tbl = Table(data, hAlign="LEFT")
            tbl.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#4B3F72")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTSIZE", (0, 0), (-1, -1), 8),
                ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#cbd5e1")),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f3f1f8")]),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ]))
            story.append(Spacer(1, 4))
            story.append(tbl)
    buf = io.BytesIO()
    SimpleDocTemplate(buf, pagesize=A4, topMargin=2 * cm, bottomMargin=2 * cm,
                      leftMargin=2 * cm, rightMargin=2 * cm).build(story)
    return buf.getvalue()


def to_pptx(payload: dict) -> bytes:
    """Genera una presentación .pptx DISEÑADA (16:9, plantilla Evalys): portada, agenda,
    diapositivas de contenido con banda de encabezado, viñetas estilizadas, pie con fuente y
    numeración, y auto-ajuste del texto. Payload:
    {titulo, subtitulo, fuente, slides:[{titulo, bullets:[..]} | {titulo, texto:".."}]}.
    Si no hay `slides`, cae a `secciones`."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
    from pptx.enum.shapes import MSO_SHAPE

    INK = RGBColor(0x0F, 0x1B, 0x2D)
    ACCENT = RGBColor(0x14, 0xB8, 0xA6)
    ACCENT2 = RGBColor(0x7C, 0x5C, 0xFF)
    CARDLINE = RGBColor(0xDD, 0xE4, 0xEC)
    PAPER = RGBColor(0xFF, 0xFF, 0xFF)
    SOFT = RGBColor(0xF5, 0xF8, 0xFB)
    SLATE = RGBColor(0x2E, 0x3F, 0x57)
    MUTED = RGBColor(0x8A, 0x99, 0xAD)
    LIGHTINK = RGBColor(0xAB, 0xBC, 0xD0)
    W, H = 13.333, 7.5
    FONT = "Calibri"

    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    def _no_border(shp):
        shp.line.fill.background()
        try:
            shp.shadow.inherit = False
        except Exception:  # noqa: BLE001
            pass

    def rect(slide, x, y, w, h, color, to_back=False):
        r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        r.fill.solid(); r.fill.fore_color.rgb = color
        _no_border(r)
        if to_back:
            sp = r._element
            sp.getparent().remove(sp)
            slide.shapes._spTree.insert(2, sp)
        return r

    def box(slide, x, y, w, h, anchor=None, fit=False):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        if anchor is not None:
            tf.vertical_anchor = anchor
        if fit:
            try:
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception:  # noqa: BLE001
                pass
        return tf

    def run(p, text, size, color, bold=False):
        r = p.add_run(); r.text = text
        f = r.font; f.size = Pt(size); f.bold = bold; f.name = FONT; f.color.rgb = color
        return r

    titulo = str(payload.get("titulo") or "Material de clase")
    subtitulo = str(payload.get("subtitulo") or "")
    fuente = str(payload.get("fuente") or "")

    slides = payload.get("slides")
    if not slides:
        slides = [{"titulo": sec.get("heading", ""), "texto": sec.get("texto", "")}
                  for sec in (payload.get("secciones") or [])]
    slides = slides or []

    # ── Portada ──
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, W, H, INK, to_back=True)
    rect(s, 0, 0, 0.32, H, ACCENT)                       # barra lateral de acento
    tf = box(s, 1.0, 0.95, 11.5, 0.5); run(tf.paragraphs[0], "MATERIAL DE CLASE · EVALYS", 13, ACCENT, True)
    rect(s, 1.03, 1.55, 1.7, 0.09, ACCENT)               # subrayado
    tf = box(s, 0.97, 2.0, 11.4, 3.3, anchor=MSO_ANCHOR.TOP, fit=True)
    run(tf.paragraphs[0], titulo, 40, PAPER, True)
    if subtitulo:
        tf = box(s, 1.0, 5.5, 11.4, 1.0, fit=True); run(tf.paragraphs[0], subtitulo, 18, LIGHTINK)
    if fuente:
        tf = box(s, 1.0, 6.85, 11.4, 0.4); run(tf.paragraphs[0], "Fuente: " + fuente[:110], 10, MUTED)

    # Agenda (roadmap) si hay suficientes secciones
    titles = [str(sl.get("titulo") or "") for sl in slides if str(sl.get("titulo") or "").strip()]
    deck = []
    if len(titles) >= 3:
        deck.append({"titulo": "En esta sesión", "bullets": titles, "_agenda": True})
    deck.extend(slides)

    def _split(sl):
        raw = sl.get("bullets")
        prose = False
        if not raw and sl.get("texto"):
            txt = str(sl["texto"])
            if "\n" in txt.strip():
                raw = [ln.strip() for ln in txt.split("\n") if ln.strip()]
            else:
                raw = [txt.strip()]; prose = True
        return [str(b) for b in (raw or []) if str(b).strip()], prose

    # Aplana el mazo en UNIDADES de diapositiva (respetando chunking y layout 'duo').
    units = []
    for sl in deck:
        title = str(sl.get("titulo") or "")
        kicker = str(sl.get("kicker") or "")
        if sl.get("layout") == "duo" and sl.get("cols"):
            units.append({"kind": "duo", "title": title, "kicker": kicker, "cols": sl["cols"]})
            continue
        items, prose = _split(sl)
        chunks = [items[i:i + 6] for i in range(0, len(items), 6)] or [[]]
        for ci, chunk in enumerate(chunks):
            units.append({"kind": "std", "title": title + ("" if ci == 0 else " (cont.)"),
                          "kicker": kicker, "items": chunk, "prose": prose})
    total = len(units)

    def header(sc, title, kicker, idx):
        rect(sc, 0, 0, W, 1.5, INK)
        rect(sc, 0, 1.5, W, 0.09, ACCENT)
        yt = 0.34
        if kicker:
            kb = box(sc, 0.9, 0.24, 10.4, 0.34)
            run(kb.paragraphs[0], kicker.upper(), 11, ACCENT, True)
            yt = 0.62
        th = box(sc, 0.9, yt, 10.5, 0.78, anchor=MSO_ANCHOR.TOP, fit=True)
        run(th.paragraphs[0], title, 22, PAPER, True)
        tn = box(sc, 11.5, 0.52, 1.5, 0.55); pn = tn.paragraphs[0]; pn.alignment = PP_ALIGN.RIGHT
        run(pn, str(idx) + " / " + str(total), 11, LIGHTINK, True)

    def card(sc, x, y, w, h, hcolor, htext, bullets):
        c = sc.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        c.fill.solid(); c.fill.fore_color.rgb = PAPER
        c.line.color.rgb = CARDLINE; c.line.width = Pt(1)
        try:
            c.shadow.inherit = False
        except Exception:  # noqa: BLE001
            pass
        hh = 0.66
        hs = sc.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(hh))
        hs.fill.solid(); hs.fill.fore_color.rgb = hcolor; _no_border(hs)
        ht = box(sc, x + 0.28, y + 0.05, w - 0.56, hh - 0.08, anchor=MSO_ANCHOR.MIDDLE, fit=True)
        run(ht.paragraphs[0], htext, 15, PAPER, True)
        bt = box(sc, x + 0.3, y + hh + 0.16, w - 0.6, h - hh - 0.34, anchor=MSO_ANCHOR.TOP, fit=True)
        first = True
        for b in bullets:
            p = bt.paragraphs[0] if first else bt.add_paragraph()
            first = False
            p.space_after = Pt(8); p.line_spacing = 1.05
            run(p, "▸  ", 13, hcolor, True)
            run(p, str(b), 13, SLATE)

    idx = 0
    for u in units:
        idx += 1
        sc = prs.slides.add_slide(blank)
        rect(sc, 0, 0, W, H, SOFT, to_back=True)
        header(sc, u["title"], u["kicker"], idx)
        if u["kind"] == "duo":
            cols = (u["cols"] or [])[:2]
            colors = [ACCENT, ACCENT2]
            xs = [0.9, 6.9]
            for k, col in enumerate(cols):
                card(sc, xs[k], 1.95, 5.55, 4.5, colors[k],
                     str(col.get("titulo") or ("Postura " + ("A" if k == 0 else "B"))),
                     [str(b) for b in (col.get("bullets") or []) if str(b).strip()])
            if len(cols) == 2:
                vb = sc.shapes.add_shape(MSO_SHAPE.OVAL, Inches(W / 2 - 0.42), Inches(3.85), Inches(0.84), Inches(0.84))
                vb.fill.solid(); vb.fill.fore_color.rgb = INK; _no_border(vb)
                vt = box(sc, W / 2 - 0.42, 3.95, 0.84, 0.64, anchor=MSO_ANCHOR.MIDDLE)
                pv = vt.paragraphs[0]; pv.alignment = PP_ALIGN.CENTER; run(pv, "VS", 16, PAPER, True)
        else:
            rect(sc, 0.9, 2.1, 0.06, 4.4, ACCENT)
            body = box(sc, 1.25, 2.0, 11.1, 4.6, anchor=MSO_ANCHOR.TOP, fit=True)
            first = True
            for b in u["items"]:
                p = body.paragraphs[0] if first else body.add_paragraph()
                first = False
                p.space_after = Pt(12); p.line_spacing = 1.08
                if u["prose"]:
                    run(p, b, 17, SLATE)
                else:
                    run(p, "▸  ", 18, ACCENT, True)
                    run(p, b, 18, SLATE)
        if fuente:
            fb = box(sc, 0.9, 6.98, 11.5, 0.35)
            run(fb.paragraphs[0], "Evalys · " + fuente[:95], 9, MUTED)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


MEDIA = {
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pdf": "application/pdf",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def exportar(formato: str, payload: dict) -> tuple[bytes, str]:
    if formato == "xlsx":
        return to_xlsx(payload), MEDIA["xlsx"]
    if formato == "docx":
        return to_docx(payload), MEDIA["docx"]
    if formato == "pdf":
        return to_pdf(payload), MEDIA["pdf"]
    if formato == "pptx":
        return to_pptx(payload), MEDIA["pptx"]
    raise ValueError("formato no soportado")
